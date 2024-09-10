from pystrich.datamatrix import DataMatrixEncoder
import pandas as pd 
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

df = pd.read_csv("tarot.csv")

ppsavename = "tarot.pptx"
prs = Presentation()

slide_width = Cm(4.4)
slide_height = Cm(8.0)

prs.slide_width = slide_width
prs.slide_height = slide_height

blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

for i in range(0, 22):
    num = str(i)
    num_reverse = str(i+22)
    nameEN = df.iloc[i,1]
    nameJA = df.iloc[i,2]
    
    print(i,nameEN)

    # generate datamatrix(正位置) with number "i"
    # number: 0-21
    encoder = DataMatrixEncoder(num)
    encoder.save(f"temp.png")

    # make slide and add image
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(
        "temp.png",
        left=Inches(0), top=Inches(0),
        width=Inches(0.5), height=Inches(0.5)
    )

    # generate datamatrix(逆位置) with number "i + 22"  
    # number: 22 - 43
    encoder = DataMatrixEncoder(num_reverse)
    encoder.save(f"tempRV.png")

    slide.shapes.add_picture(
        "tempRV.png",
        left=Inches(1.225), top=Inches(2.65),
        width=Inches(0.5), height=Inches(0.5)
    )


    # add text
    title = slide.shapes.add_textbox(
        left=Cm(0), top=Cm(0),
        width=slide_width-1, height=Cm(2)
    )

    text_frame = title.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # テキストボックス全体を中央揃えに設定

    p = text_frame.add_paragraph()
    p.text = (str(i) + " "  + str(nameEN))
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.bold = True

    title = slide.shapes.add_textbox(left=Cm(0), top=Cm(0.65),width=slide_width-1, height=Cm(2))
    text_frame = title.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # テキストボックス全体を中央揃えに設定

    p = text_frame.add_paragraph()
    p.text = nameJA
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.bold = True

# save slide
prs.save(ppsavename)